import os
from typing import List, Dict, Any
import logging
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import pandas as pd
from PIL import Image
import io
import pypdf
import pytesseract
import easyocr
import numpy as np

logger = logging.getLogger(__name__)

class UploadService:
    @staticmethod
    def parse_docx(file_content: bytes) -> str:
        doc = DocxDocument(io.BytesIO(file_content))
        return "\n".join([para.text for para in doc.paragraphs])

    @staticmethod
    def parse_pptx(file_content: bytes) -> str:
        prs = Presentation(io.BytesIO(file_content))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    @staticmethod
    def parse_excel(file_content: bytes) -> str:
        # Using pandas for easier excel to text conversion
        df_dict = pd.read_excel(io.BytesIO(file_content), sheet_name=None)
        output = []
        for sheet_name, df in df_dict.items():
            output.append(f"Sheet: {sheet_name}")
            output.append(df.to_string())
        return "\n\n".join(output)

    @staticmethod
    def parse_pdf(file_content: bytes) -> str:
        reader = pypdf.PdfReader(io.BytesIO(file_content))
        text = []
        for page in reader.pages:
            content = page.extract_text()
            if content:
                text.append(content)
        return "\n".join(text)

    @staticmethod
    def parse_image(file_content: bytes) -> str:
        try:
            img = Image.open(io.BytesIO(file_content))
            
            # Diagnostic: Check if Tesseract is in PATH
            tesseract_available = False
            try:
                pytesseract.get_tesseract_version()
                tesseract_available = True
            except Exception:
                logger.warning("Pytesseract: Tesseract-OCR binary not found in system path.")

            # Tier 1: Pytesseract
            if tesseract_available:
                try:
                    text = pytesseract.image_to_string(img)
                    if text.strip():
                        return text
                except Exception as e:
                    logger.debug(f"Pytesseract execution failed: {e}")

            # Tier 2: EasyOCR
            try:
                # Optimized: Only init reader once if possible, but here we stay simple for stability
                reader = easyocr.Reader(['en'], gpu=False) # GPU False for broader compatibility
                results = reader.readtext(np.array(img))
                text = " ".join([res[1] for res in results])
                if text.strip():
                    return text
            except Exception as e:
                logger.debug(f"EasyOCR failed: {e}")

            # Fallback: Metadata + Warning
            msg = f"Image file processed ({img.format}, {img.size})."
            if not tesseract_available:
                msg += " NOTE: Tesseract-OCR is not installed on this server; results may be limited."
            return f"{msg} (No clear text extracted via OCR)"
        except Exception as e:
            logger.error(f"Image parsing failed: {e}")
            return f"Error: Failed to process image file. {str(e)}"

    def parse_file(self, file_content: bytes, filename: str) -> str:
        ext = os.path.splitext(filename)[1].lower()
        try:
            if ext in ['.docx']:
                return self.parse_docx(file_content)
            elif ext in ['.pptx']:
                return self.parse_pptx(file_content)
            elif ext in ['.xlsx', '.xls']:
                return self.parse_excel(file_content)
            elif ext in ['.png', '.jpg', '.jpeg', '.webp']:
                return self.parse_image(file_content)
            elif ext in ['.pdf']:
                return self.parse_pdf(file_content)
            elif ext in ['.txt', '.md', '.csv']:
                return file_content.decode('utf-8', errors='ignore')
            else:
                return f"Unsupported file format: {ext}"
        except Exception as e:
            logger.error(f"Failed to parse {filename}: {e}")
            return f"Error parsing {filename}: {str(e)}"
