import pdfplumber
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
import openpyxl
import io, os, re
import base64
from typing import Dict, Any, List

class FileProcessor:
    """
    Handles extraction of text and metadata from various file formats.
    Supports: PDF, DOCX, XLSX, CSV, PPTX, Images (OCR), TXT/MD.
    """
    
    @staticmethod
    async def process_file(file_bytes: bytes, filename: str) -> Dict[str, Any]:
        """
        Master method to route file processing based on extension.
        """
        ext = filename.lower().split('.')[-1]
        
        processors = {
            'pdf': FileProcessor.process_pdf,
            'png': FileProcessor.process_image,
            'jpg': FileProcessor.process_image,
            'jpeg': FileProcessor.process_image,
            'webp': FileProcessor.process_image,
            'gif': FileProcessor.process_image,
            'docx': FileProcessor.process_docx,
            'doc': FileProcessor.process_docx,
            'xlsx': FileProcessor.process_excel,
            'xls': FileProcessor.process_excel,
            'csv': FileProcessor.process_csv,
            'pptx': FileProcessor.process_pptx,
            'txt': FileProcessor.process_text,
            'md': FileProcessor.process_text,
            'py': FileProcessor.process_text,
            'js': FileProcessor.process_text,
            'html': FileProcessor.process_text,
            'css': FileProcessor.process_text,
            'json': FileProcessor.process_text,
        }

        processor = processors.get(ext)
        if not processor:
            return {
                "success": False,
                "error": f"File type .{ext} not supported",
                "supported": list(processors.keys())
            }

        try:
            result = await processor(file_bytes, filename)
            result['filename'] = filename
            result['extension'] = ext
            
            # Semantic chunking for RAG
            if result.get('text'):
                result['chunks'] = FileProcessor.chunk_text(result['text'])
                
            return result
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "filename": filename
            }

    @staticmethod
    async def process_pdf(file_bytes: bytes, filename: str) -> Dict[str, Any]:
        text_pages = []
        total_text = ""
        
        try:
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                total_pages = len(pdf.pages)
                
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    
                    if page_text and len(page_text.strip()) > 50:
                        text_pages.append({
                            "page": i + 1,
                            "text": page_text.strip(),
                            "method": "text_layer"
                        })
                        total_text += f"\n[Page {i+1}]\n{page_text.strip()}\n"
                    else:
                        # Fallback to OCR if installed
                        try:
                            # Note: Tesseract must be installed on system for this to work
                           ocr_text = "" # Placeholder if tesseract fails/not installed
                           # ocr_text = pytesseract.image_to_string(page.to_image(resolution=150).original)
                           if ocr_text.strip():
                               text_pages.append({"page": i+1, "text": ocr_text, "method": "ocr"})
                               total_text += f"\n[Page {i+1} - OCR]\n{ocr_text.strip()}\n"
                        except:
                             pass # OCR failed or not available

            return {
                "success": True,
                "type": "pdf",
                "text": total_text,
                "pages": text_pages,
                "page_count": total_pages
            }
        except Exception as e:
            return {"success": False, "error": f"PDF Error: {str(e)}"}

    @staticmethod
    async def process_docx(file_bytes: bytes, filename: str) -> Dict[str, Any]:
        doc = Document(io.BytesIO(file_bytes))
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                if 'Heading' in para.style.name:
                    text_parts.append(f"\n## {para.text.strip()}\n")
                else:
                    text_parts.append(para.text.strip())
                    
        for i, table in enumerate(doc.tables):
            text_parts.append(f"\n[Table {i+1}]")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                text_parts.append(" | ".join(cells))
                
        return {
            "success": True,
            "type": "docx",
            "text": "\n".join(text_parts),
            "paragraph_count": len(doc.paragraphs)
        }

    @staticmethod
    async def process_excel(file_bytes: bytes, filename: str) -> Dict[str, Any]:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        all_text = ""
        sheets = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_text += f"\n[Sheet: {sheet_name}]\n"
            rows = []
            for row in ws.iter_rows(values_only=True):
                # Filter None values
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in row_data):
                    rows.append(" | ".join(row_data))
            
            sheet_text = "\n".join(rows)
            all_text += sheet_text + "\n"
            sheets.append({"name": sheet_name, "rows": len(rows)})
            
        return {
            "success": True,
            "type": "excel",
            "text": all_text,
            "sheets": sheets
        }

    @staticmethod
    async def process_csv(file_bytes: bytes, filename: str) -> Dict[str, Any]:
        import csv
        text = file_bytes.decode('utf-8', errors='replace')
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        
        formatted = ""
        if rows:
            header_row = rows[0]
            formatted += "Headers: " + " | ".join(header_row) + "\n\n"
            data_rows = rows[1:100] # Limit to first 100 rows for context
            for row in data_rows:
                formatted += " | ".join(row) + "\n"
            if len(rows) > 100:
                formatted += f"\n[{len(rows)-100} more rows omitted]"
                
        return {
            "success": True,
            "type": "csv",
            "text": formatted,
            "row_count": len(rows)
        }
        
    @staticmethod
    async def process_pptx(file_bytes: bytes, filename: str) -> Dict[str, Any]:
        prs = Presentation(io.BytesIO(file_bytes))
        full_text = ""
        slides = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                # Use getattr and check if it's a string to satisfy strict type checkers
                text_val = getattr(shape, "text", None)
                if isinstance(text_val, str) and text_val.strip():
                    slide_text.append(text_val)
            
            joined = "\n".join(slide_text)
            if joined.strip():
                full_text += f"\n[Slide {i+1}]\n{joined}\n"
                slides.append({"slide": i+1, "text": joined})
                
        return {
            "success": True,
            "type": "pptx",
            "text": full_text,
            "slide_count": len(slides)
        }

    @staticmethod
    async def process_text(file_bytes: bytes, filename: str) -> Dict[str, Any]:
        text = file_bytes.decode('utf-8', errors='replace')
        return {
            "success": True,
            "type": "text",
            "text": text,
            "char_count": len(text)
        }

    @staticmethod
    async def process_image(file_bytes: bytes, filename: str) -> Dict[str, Any]:
        """
        Processes images. Uses Vision model by default, but provides local OCR fallback.
        """
        text = ""
        try:
             image = Image.open(io.BytesIO(file_bytes))
             # Check if tesseract is available before running
             try:
                 pytesseract.get_tesseract_version()
                 text = pytesseract.image_to_string(image)
             except:
                 # Local OCR not available (binary missing)
                 text = "[Image analysis will be performed by Vision Model]"
        except Exception as e:
             text = f"[Image read error: {str(e)}]"
             
        # Return base64 for frontend display or LLM vision
        img_b64 = base64.b64encode(file_bytes).decode('utf-8')
        
        return {
            "success": True,
            "type": "image",
            "text": text,
            "image_base64": img_b64,
            "vision_ready": True
        }

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> List[Dict[str, Any]]:
        """
        Splits text into overlapping chunks for RAG.
        """
        if not text: return []
        
        words = text.split()
        chunks = []
        i = 0
        
        while i < len(words):
            chunk_words = words[i:i + chunk_size]
            chunk_text = " ".join(chunk_words)
            chunks.append({
                "index": len(chunks),
                "text": chunk_text,
                "word_count": len(chunk_words)
            })
            i += (chunk_size - overlap)
            
        return chunks
